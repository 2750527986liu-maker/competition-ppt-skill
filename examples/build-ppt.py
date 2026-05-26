#!/usr/bin/env python3
"""
助你好梦 - 中国国际大学生创新大赛 路演PPT
使用原创设计素材，学习国赛模板设计思路
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, copy

# === CONFIG ===
W = Inches(13.33)
H = Inches(7.5)
IMG = "/sessions/zealous-eager-fermat/mnt/outputs/ppt-images"
BG = "/sessions/zealous-eager-fermat/mnt/被子/ppt-assets"
OUT = "/sessions/zealous-eager-fermat/mnt/被子/助你好梦_国赛路演.pptx"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# Colors
NAVY = RGBColor(0x0A, 0x16, 0x28)
DARK_BLUE = RGBColor(0x0F, 0x1D, 0x35)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
LIGHT_GOLD = RGBColor(0xE8, 0xD5, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
SUB_TEXT = RGBColor(0x64, 0x74, 0x8B)
ACCENT_BLUE = RGBColor(0x38, 0x6F, 0xC4)
SOFT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank

# === HELPERS ===
def add_bg(slide, bg_path):
    """Add background image to slide"""
    if os.path.exists(bg_path):
        slide.background.fill.background()
        pic = slide.shapes.add_picture(bg_path, Emu(0), Emu(0), W, H)
        # Move to back
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

def text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT,
             bold=False, align=PP_ALIGN.LEFT, font_name="等线", anchor=MSO_ANCHOR.TOP,
             line_spacing=1.5):
    """Add a text box with specified formatting"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    if anchor:
        tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def multi_text(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT,
               bold=False, align=PP_ALIGN.LEFT, font_name="等线", line_spacing=1.5):
    """Add text box with multiple paragraphs"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def gold_line(slide, left, top, width):
    """Add a gold horizontal accent line"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def section_badge(slide, text, left=0.3, top=0.2):
    """Add section label badge (top-left, like template)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    shape.rotation = 0
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(22)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "等线"
    return shape

def sub_header(slide, text, left=3.2, top=0.2):
    """Add subtitle with —— prefix (like template)"""
    return text_box(slide, left, top, 9.5, 0.6, f"——{text}", font_size=16, color=SUB_TEXT)

def guosai_icon(slide, left=12.2, top=0.1):
    """Add the 国赛 icon to top right"""
    icon_path = os.path.join(BG, "guosai-icon.png")
    if os.path.exists(icon_path):
        slide.shapes.add_picture(icon_path, Inches(left), Inches(top), Inches(0.9), Inches(0.9))

def card_shape(slide, left, top, width, height, fill_color=WHITE, alpha=None):
    """Add a card (rounded rectangle) with optional transparency"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    return shape

def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if file exists"""
    if os.path.exists(path):
        if width and height:
            return slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            return slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
        else:
            return slide.shapes.add_picture(path, Inches(left), Inches(top))

def add_page_number(slide, num, total=20):
    """Add page number bottom right"""
    text_box(slide, 12.0, 6.9, 1.0, 0.4, f"{num}/{total}", font_size=10, color=SUB_TEXT, align=PP_ALIGN.RIGHT)

# ========================================
# SLIDE 1: COVER
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "cover-bg.png"))

# Gold accent decoration
add_image_safe(s, os.path.join(BG, "gold-line.png"), 0.8, 2.35, width=5)

# Moon icon
add_image_safe(s, os.path.join(BG, "moon-icon.png"), 10.5, 1.0, width=1.2)

# Competition title top left
text_box(s, 0.8, 0.2, 6, 0.5, "中国国际大学生创新大赛（2025）", font_size=16, color=LIGHT_GOLD, bold=True)

# Main title
text_box(s, 0.8, 2.5, 7, 1.2, "助你好梦", font_size=56, color=WHITE, bold=True)

# Subtitle
text_box(s, 0.8, 3.8, 7, 0.6, "重新定义现代压力下的深度睡眠方案", font_size=22, color=LIGHT_GOLD)

# Gold separator
gold_line(s, 0.8, 4.6, 5)

# Team info cards
for i, (label, value) in enumerate([("团队", "助你好梦创业团队"), ("赛道", "高教主赛道"), ("类别", "产品创新")]):
    x = 0.8 + i * 2.5
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.2), Inches(2.2), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x40)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{label}  {value}"
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE
    run.font.name = "等线"

# ========================================
# SLIDE 2: SLEEP CRISIS - DATA
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "睡眠危机")
sub_header(s, "当代大学生与中青年群体的睡眠质量现状")
# (guosai icon removed - using original design)

# Main data cards
data_items = [
    ("3亿+", "中国睡眠障碍人群", "我国存在睡眠障碍的人口规模"),
    ("46.9%", "大学生失眠比例", "基于META综述的循证数据"),
    ("7.2h", "推荐最低睡眠时长", "而大量学生远低于此标准"),
]
for i, (num, title, desc) in enumerate(data_items):
    x = 0.5 + i * 4.2
    card_shape(s, x, 1.5, 3.8, 2.8)
    text_box(s, x + 0.3, 1.7, 3.2, 0.8, num, font_size=48, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.3, 2.6, 3.2, 0.5, title, font_size=18, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.3, 3.2, 3.2, 0.8, desc, font_size=14, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# Literature images in cards
card_shape(s, 0.5, 4.5, 3.8, 2.5)
add_image_safe(s, os.path.join(IMG, "46.9%的大学生失眠META综述.jpg"), 0.7, 4.6, width=3.4)
text_box(s, 0.7, 6.6, 3.4, 0.3, "META综述：46.9%大学生失眠", font_size=10, color=SUB_TEXT, align=PP_ALIGN.CENTER)

card_shape(s, 4.6, 4.5, 3.8, 2.5)
add_image_safe(s, os.path.join(IMG, "香港年轻人 Journal of Psychiatric Research.png"), 4.8, 4.6, width=3.4)
text_box(s, 4.8, 6.6, 3.4, 0.3, "香港年轻人睡眠研究", font_size=10, color=SUB_TEXT, align=PP_ALIGN.CENTER)

card_shape(s, 8.7, 4.5, 3.8, 2.5)
add_image_safe(s, os.path.join(IMG, "加州年轻人Cureus.png"), 8.9, 4.6, width=3.4)
text_box(s, 8.9, 6.6, 3.4, 0.3, "加州年轻人睡眠研究", font_size=10, color=SUB_TEXT, align=PP_ALIGN.CENTER)

add_page_number(s, 2)

# ========================================
# SLIDE 3: SCIENTIFIC BASIS
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "科学依据")
sub_header(s, "深层压力刺激——被循证医学验证的助眠机制")
# (guosai icon removed - using original design)

# Key point
text_box(s, 0.5, 1.3, 12.3, 0.6,
         "深层压力刺激（Deep Pressure Stimulation）能够显著降低皮质醇水平，提高褪黑素分泌，促进入睡",
         font_size=16, color=DARK_TEXT, bold=True)

gold_line(s, 0.5, 1.95, 12.3)

# Literature cards - 3 columns
lit_data = [
    ("压力与褪黑素", "深层压力刺激通过激活副交感神经系统，促进褪黑素分泌，缩短入睡潜伏期", os.path.join(IMG, "薰衣草 Experimental Gerontology .png")),
    ("微气候影响", "被窝温度每升高1°C，入睡时间显著延长。传统加重毯导致温度过高", os.path.join(IMG, "微气候 Journal of Thermal Biology.png")),
    ("温度与睡眠", "Nature子刊研究证实：适宜的热环境是良好睡眠的基础条件", os.path.join(IMG, "温度Nature子刊.png")),
]
for i, (title, desc, img_path) in enumerate(lit_data):
    x = 0.5 + i * 4.2
    card_shape(s, x, 2.2, 3.8, 4.8)
    add_image_safe(s, img_path, x + 0.2, 2.4, width=3.4, height=2.0)
    text_box(s, x + 0.2, 4.5, 3.4, 0.4, title, font_size=16, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    gold_line(s, x + 1.0, 5.0, 1.8)
    text_box(s, x + 0.2, 5.1, 3.4, 1.5, desc, font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_page_number(s, 3)

# ========================================
# SLIDE 4: INDUSTRY PAIN POINTS
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "行业痛点")
sub_header(s, "传统加重毯的三座大山——为什么实验室有效，商业却失败")
# (guosai icon removed - using original design)

pain_points = [
    ("难清洗", "传统加重毯配重层不可拆卸，无法机洗，卫生问题严重。用户购买后实际使用率低。", "❌"),
    ("太闷热", "被窝微气候温度过高。加重毯增加的重量同时大幅降低透气性，导致体温升高，入睡反而更困难。", "🌡"),
    ("不可调", "配重固定无法调节，无法适应不同身体部位的需求差异。每个人的肩、腹、腿对压力需求不同。", "🔒"),
]
for i, (title, desc, icon) in enumerate(pain_points):
    x = 0.5 + i * 4.2
    card_shape(s, x, 1.5, 3.8, 3.5)
    # Icon circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.4), Inches(1.7), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xEF, 0x44, 0x44) if i == 0 else (RGBColor(0xF5, 0x9E, 0x0B) if i == 1 else RGBColor(0x6B, 0x72, 0x80))
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.size = Pt(28)
    text_box(s, x + 0.3, 2.8, 3.2, 0.4, title, font_size=22, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.3, 3.3, 3.2, 1.5, desc, font_size=13, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# Word cloud in card
card_shape(s, 2.0, 5.1, 9.0, 1.7)
add_image_safe(s, os.path.join(IMG, "词云2.jpg"), 2.2, 5.2, width=4.2, height=1.5)
add_image_safe(s, os.path.join(IMG, "wordcloud.png"), 6.6, 5.2, width=4.2, height=1.5)

# Bottom insight
text_box(s, 0.5, 6.6, 12.3, 0.6, "为什么加重毯在实验室有效，商业实用中却退货率居高不下？罪魁祸首是被窝微气候。",
         font_size=14, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

add_page_number(s, 4)

# ========================================
# SLIDE 5: PRODUCT INTRODUCTION
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "产品介绍")
sub_header(s, "模块化外挂式加重助眠被套系统——第二代创新产品")
# (guosai icon removed - using original design)

# Left: product description
card_shape(s, 0.5, 1.5, 6.0, 5.3)
text_box(s, 0.8, 1.7, 5.5, 0.5, "助你好梦·加重被套", font_size=28, color=NAVY, bold=True)
gold_line(s, 0.8, 2.3, 3)

features = [
    "一款可自由调节配重位置与重量的模块化被套系统",
    "根据身体不同区域（肩、腹、腿）的敏感度，精准粘贴加重模块",
    "像贴魔术贴一样操作，实现精准深压触感而不过多影响空气流通",
    "外挂式设计：被套可正常机洗，重量不挤压内部被芯",
    "彻底解决传统加重毯清洗难、闷热、功能单一的痛点",
]
multi_text(s, 0.8, 2.5, 5.5, 4.0, features, font_size=15, color=DARK_TEXT, line_spacing=1.8)

# Right: product images in cards
card_shape(s, 6.8, 1.3, 5.8, 2.8)
add_image_safe(s, os.path.join(IMG, "图1.png"), 7.0, 1.4, width=5.4, height=2.5)
card_shape(s, 6.8, 4.3, 5.8, 2.8)
add_image_safe(s, os.path.join(IMG, "图2.png"), 7.0, 4.4, width=5.4, height=2.5)

add_page_number(s, 5)

# ========================================
# SLIDE 6: TECHNICAL INNOVATION
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "技术创新")
sub_header(s, "模块化设计——被套与配重包分离化")
# (guosai icon removed - using original design)

# Innovation cards in 2x2 grid
innovations = [
    ("模块化分区", "五区独立配重系统：左/右胸部、腹部中央、左/右腿部，互不干扰，真正个性化配置", "📐"),
    ("双通道放松", "每个模块填充不锈钢丸与天然中草药（薰衣草、柠檬草），物理重力+嗅觉双通道", "🌿"),
    ("透气不闷热", "外挂式设计不覆盖被芯表面，保持空气流通。不像传统加重毯覆盖整床导致闷热", "💨"),
    ("便捷清洗", "模块通过魔术贴固定，可单独拆卸。被套正常机洗，配重包擦拭即可，方便卫生", "🧺"),
]
for i, (title, desc, icon) in enumerate(innovations):
    row, col = divmod(i, 2)
    x = 0.5 + col * 6.3
    y = 1.3 + row * 3.0
    card_shape(s, x, y, 5.9, 2.6)
    text_box(s, x + 0.3, y + 0.2, 5.3, 0.5, f"{icon}  {title}", font_size=22, color=NAVY, bold=True)
    gold_line(s, x + 0.3, y + 0.8, 2.5)
    text_box(s, x + 0.3, y + 1.0, 5.3, 1.4, desc, font_size=14, color=DARK_TEXT, line_spacing=1.6)

add_page_number(s, 6)

# ========================================
# SLIDE 7: MICROCLIMATE LITERATURE
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "文献支撑")
sub_header(s, "微气候与温度——来自Nature子刊与顶刊的循证研究")
# (guosai icon removed - using original design)

# Literature images in 3 columns
papers = [
    ("微气候研究", "Journal of Thermal Biology", os.path.join(IMG, "微气候 Journal of Thermal Biology.png")),
    ("温度与睡眠", "Nature子刊", os.path.join(IMG, "温度Nature子刊.png")),
    ("META综述", "46.9%大学生失眠", os.path.join(IMG, "46.9%的大学生失眠META综述.jpg")),
]
for i, (title, source, img_path) in enumerate(papers):
    x = 0.5 + i * 4.2
    card_shape(s, x, 1.3, 3.8, 5.5)
    add_image_safe(s, img_path, x + 0.2, 1.5, width=3.4, height=3.5)
    text_box(s, x + 0.2, 5.2, 3.4, 0.4, title, font_size=16, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    gold_line(s, x + 1.0, 5.7, 1.8)
    text_box(s, x + 0.2, 5.8, 3.4, 0.4, source, font_size=12, color=SUB_TEXT, align=PP_ALIGN.CENTER)

add_page_number(s, 7)

# ========================================
# SLIDE 8: ADHD & SPECIAL POPULATIONS
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "目标人群")
sub_header(s, "微失眠焦虑人群与ADHD儿童——深层压力的双重价值")
# (guosai icon removed - using original design)

# Left card
card_shape(s, 0.5, 1.3, 5.9, 5.5)
text_box(s, 0.8, 1.5, 5.3, 0.5, "成人失眠与焦虑", font_size=22, color=NAVY, bold=True)
gold_line(s, 0.8, 2.1, 2.5)
multi_text(s, 0.8, 2.3, 5.3, 4.0, [
    "微失眠、焦虑人群：褪黑素分泌不足，交感神经过度兴奋",
    "深层压力刺激通过激活副交感神经，自然促进褪黑素分泌",
    "替代药物干预，无副作用，安全高效",
    "覆盖16-32岁上升期青年群体",
], font_size=14, color=DARK_TEXT, line_spacing=1.8)

# Right card
card_shape(s, 6.8, 1.3, 5.9, 5.5)
text_box(s, 7.1, 1.5, 5.3, 0.5, "ADHD儿童辅助", font_size=22, color=NAVY, bold=True)
gold_line(s, 7.1, 2.1, 2.5)
multi_text(s, 7.1, 2.3, 5.3, 4.0, [
    "ADHD（注意缺陷多动障碍）儿童常伴有感觉统合失调",
    "加重毯/加重背心已在OT（作业治疗）中广泛使用",
    "深层压力帮助ADHD儿童降低焦虑、提升专注力",
    "安全性高，适合作为非药物辅助治疗手段",
], font_size=14, color=DARK_TEXT, line_spacing=1.8)

# Reference images - smaller, centered in cards
add_image_safe(s, os.path.join(IMG, "香港年轻人 Journal of Psychiatric Research.png"), 1.0, 5.4, width=4.8, height=1.2)
add_image_safe(s, os.path.join(IMG, "加州年轻人Cureus.png"), 7.3, 5.4, width=4.8, height=1.2)

add_page_number(s, 8)

# ========================================
# SLIDE 9: HERBAL INGREDIENTS
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "核心配方")
sub_header(s, "天然中草药——压力+嗅觉双通道放松系统")
# (guosai icon removed - using original design)

herbs = [
    ("薰衣草", "Lavandula angustifolia", "镇静安神，降低心率和血压。实验研究证实可显著改善睡眠质量，缩短入睡时间。", os.path.join(IMG, "薰衣草 Experimental Gerontology .png")),
    ("柠檬草", "Cymbopogon citratus", "驱蚊抗菌，清新提神又不过度兴奋。具有温和的镇静效果，适合夏季使用。", os.path.join(IMG, "柠檬草 Parasitology International.png")),
    ("柠檬香蜂草", "Melissa officinalis", "抗焦虑，促进GABA分泌。多项临床研究证实对焦虑和失眠有显著改善作用。", os.path.join(IMG, "柠檬香蜂草 Nutrients.png")),
]
for i, (name, latin, desc, img_path) in enumerate(herbs):
    x = 0.5 + i * 4.2
    card_shape(s, x, 1.3, 3.8, 5.7)
    add_image_safe(s, img_path, x + 0.2, 1.5, width=3.4, height=2.8)
    text_box(s, x + 0.2, 4.4, 3.4, 0.4, name, font_size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.2, 4.9, 3.4, 0.3, latin, font_size=11, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    gold_line(s, x + 1.0, 5.3, 1.8)
    text_box(s, x + 0.2, 5.4, 3.4, 1.3, desc, font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_page_number(s, 9)

# ========================================
# SLIDE 10: MARKET SIZE
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "市场空间")
sub_header(s, "千亿级睡眠健康赛道，蓝海机遇待挖掘")
# (guosai icon removed - using original design)

# Market chart image
add_image_safe(s, os.path.join(IMG, "市场图表.png"), 0.5, 1.5, width=7.0)

# Key market data cards
market_data = [
    ("4000亿+", "中国睡眠经济市场规模", "2024年预计突破4000亿元"),
    ("38.2%", "睡眠健康产品年增长率", "远高于传统寝具行业增速"),
    ("3亿+", "潜在用户群体", "睡眠障碍人群持续扩大"),
]
for i, (num, title, desc) in enumerate(market_data):
    y = 1.5 + i * 1.8
    card_shape(s, 8.0, y, 4.8, 1.5)
    text_box(s, 8.3, y + 0.1, 2.0, 0.6, num, font_size=32, color=ACCENT_BLUE, bold=True)
    text_box(s, 8.3, y + 0.7, 4.2, 0.3, title, font_size=14, color=DARK_TEXT, bold=True)
    text_box(s, 8.3, y + 1.05, 4.2, 0.3, desc, font_size=12, color=SUB_TEXT)

text_box(s, 0.5, 6.5, 12.3, 0.5, "数据来源：艾媒咨询、中商产业研究院、中国睡眠研究会",
         font_size=11, color=SUB_TEXT, align=PP_ALIGN.RIGHT)

add_page_number(s, 10)

# ========================================
# SLIDE 11: USER PERSONA
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "用户画像")
sub_header(s, "16-32岁青年群体——上升期中的睡眠健康需求")
# (guosai icon removed - using original design)

# Persona card
card_shape(s, 0.5, 1.3, 12.3, 5.5)

# Target user description
text_box(s, 1.0, 1.6, 11, 0.6, "核心目标用户", font_size=26, color=NAVY, bold=True)
gold_line(s, 1.0, 2.3, 3)

persona_items = [
    "年龄：16-32岁，处于学业/职业上升期的青年群体",
    "状态：轻度焦虑、入睡困难、睡眠质量不佳",
    "场景：在校大学生、考研/考公备考族、职场新人、创业者",
    "需求：高效睡眠、身体恢复、压力释放、非药物解决方案",
    "特征：接受新事物，有一定消费能力，注重生活品质",
    "典型代表：在校大学生——面临学业压力、作息不规律、宿舍环境影响",
]
multi_text(s, 1.0, 2.5, 5.5, 4.0, persona_items, font_size=15, color=DARK_TEXT, line_spacing=1.7)

# Right side: user data images
add_image_safe(s, os.path.join(IMG, "46.9%的大学生失眠META综述.jpg"), 7.0, 2.5, width=5.5)

# Bottom insight
text_box(s, 1.0, 5.5, 11, 0.6,
         "中国市场缺乏针对青年人群的专业助眠物理产品，这是巨大的商业蓝海",
         font_size=16, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

add_page_number(s, 11)

# ========================================
# SLIDE 12: BUSINESS MODEL
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "商业模式")
sub_header(s, "剃刀与刀片模式——从低频寝具到高频健康服务")
# (guosai icon removed - using original design)

# Business model explanation
card_shape(s, 0.5, 1.3, 6.0, 2.5)
text_box(s, 0.8, 1.5, 5.4, 0.5, "盈利逻辑", font_size=22, color=NAVY, bold=True)
gold_line(s, 0.8, 2.1, 2.5)
multi_text(s, 0.8, 2.3, 5.4, 1.2, [
    "被套（平台）+ 配重模块（耗材）= 剃刀与刀片模式",
    "高频次的加重模块订阅为主要赢利点",
    "每个配重模块使用寿命约18天，形成持续复购",
], font_size=14, color=DARK_TEXT, line_spacing=1.6)

# Product pricing
pricing = [
    ("被套", "129元", "基础平台，一次购买"),
    ("标准模块", "29元/组", "不锈钢丸配重，18天更换"),
    ("中药模块", "39元/组", "草本配方，提升附加值"),
    ("套装A", "249元", "被套+标准模块套装"),
    ("套装B", "299元", "被套+中药模块套装"),
]
for i, (name, price, desc) in enumerate(pricing):
    x = 7.0
    y = 1.3 + i * 1.0
    card_shape(s, x, y, 5.8, 0.85)
    text_box(s, 7.3, y + 0.05, 1.5, 0.4, name, font_size=16, color=NAVY, bold=True)
    text_box(s, 8.8, y + 0.05, 1.5, 0.4, price, font_size=18, color=ACCENT_BLUE, bold=True)
    text_box(s, 10.3, y + 0.05, 2.2, 0.4, desc, font_size=12, color=SUB_TEXT)

# Profit curve
text_box(s, 0.5, 4.2, 12.3, 0.5, "盈利增长曲线", font_size=18, color=NAVY, bold=True)
gold_line(s, 0.5, 4.7, 3)

# Growth timeline
growth = [
    ("第1年", "基础搭建\n目标30万营收\n淘宝+抖音双渠道启动"),
    ("第2年", "快速增长\n目标50万营收\n订阅模式成熟，复购率提升"),
    ("第3年", "规模扩张\n目标120万营收\n首轮融资，拓展线下渠道"),
    ("第5年", "生态布局\n全感官助眠生态\n白噪音枕头+整体设计"),
]
for i, (year, desc) in enumerate(growth):
    x = 0.5 + i * 3.15
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.9), Inches(2.9), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY if i == 3 else WHITE
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = year
    run.font.size = Pt(18)
    run.font.color.rgb = GOLD if i == 3 else NAVY
    run.font.bold = True
    run.font.name = "等线"
    for line in desc.split('\n'):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE if i == 3 else DARK_TEXT
        r.font.name = "微软雅黑"

add_page_number(s, 12)

# ========================================
# SLIDE 13: COMPETITIVE ADVANTAGES
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "竞争优势")
sub_header(s, "比药品更安全 · 比电器更实惠 · 比同类更专业")
# (guosai icon removed - using original design)

# Comparison table header
headers = ["维度", "药物助眠", "智能助眠设备", "传统加重毯", "助你好梦"]
colors_h = [NAVY, RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0x9E,0x0B), SUB_TEXT, ACCENT_BLUE]
for i, (h, c) in enumerate(zip(headers, colors_h)):
    x = 0.5 + i * 2.5
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(2.3), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = h
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "等线"

# Comparison rows
rows = [
    ("安全性", "有副作用\n依赖性风险", "电磁辐射争议", "闷热、卫生", "天然材料\n物理助眠"),
    ("有效性", "短期有效\n治标不治本", "效果因人而异", "温度过高\n适得其反", "压力+嗅觉\n双通道"),
    ("价格", "长期花费高\n月均数百元", "数千元起步", "200-800元", "129-299元\n模块持续复购"),
    ("便捷性", "需要处方", "需要充电\n操作复杂", "无法清洗", "可机洗\n即贴即用"),
]
for ri, (dim, *vals) in enumerate(rows):
    y = 2.0 + ri * 1.2
    all_vals = [dim] + list(vals)
    for ci, val in enumerate(all_vals):
        x = 0.5 + ci * 2.5
        card_shape(s, x, y, 2.3, 1.05, fill_color=WHITE)
        text_box(s, x + 0.1, y + 0.05, 2.1, 0.95, val, font_size=11, color=DARK_TEXT if ci != 4 else ACCENT_BLUE,
                bold=(ci==0 or ci==4), align=PP_ALIGN.CENTER)

add_page_number(s, 13)

# ========================================
# SLIDE 14: SALES PERFORMANCE
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "销售成果")
sub_header(s, "淘宝+抖音双渠道运营，初战告捷")
# (guosai icon removed - using original design)

# Taobao card
card_shape(s, 0.5, 1.3, 5.9, 5.5)
text_box(s, 0.8, 1.5, 5.3, 0.5, "淘宝平台", font_size=24, color=NAVY, bold=True)
gold_line(s, 0.8, 2.1, 2.5)
multi_text(s, 0.8, 2.3, 5.3, 4.0, [
    "店铺已正式上线运营",
    "SKU覆盖被套、标准模块、中药模块、套装",
    "月均销量稳步增长",
    "用户好评率高，复购意愿强",
    "毛利率达到行业领先水平",
], font_size=15, color=DARK_TEXT, line_spacing=1.8)

# Douyin card
card_shape(s, 6.8, 1.3, 5.9, 5.5)
text_box(s, 7.1, 1.5, 5.3, 0.5, "抖音直播", font_size=24, color=NAVY, bold=True)
gold_line(s, 7.1, 2.1, 2.5)
multi_text(s, 7.1, 2.3, 5.3, 4.0, [
    "抖音直播带货已启动",
    "短视频种草+直播转化模式",
    "目标用户精准触达",
    "直播数据表现优异",
    "已建立稳定的直播团队",
], font_size=15, color=DARK_TEXT, line_spacing=1.8)

add_page_number(s, 14)

# ========================================
# SLIDE 15: FINANCIAL ANALYSIS
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "财务分析")
sub_header(s, "三年财务规划——预计第二年实现盈亏平衡")
# (guosai icon removed - using original design)

# Financial data table
fin_headers = ["指标", "第一年", "第二年", "第三年"]
fin_data = [
    ("营业收入", "30万", "50万", "120万"),
    ("净利润", "-8万", "25万", "66万"),
    ("毛利率", "65%", "68%", "72%"),
    ("融资计划", "—", "—", "首轮20万/10%"),
]

# Header row
for i, h in enumerate(fin_headers):
    x = 1.0 + i * 2.8
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(2.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = h
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "等线"

# Data rows
for ri, row in enumerate(fin_data):
    y = 2.2 + ri * 0.9
    for ci, val in enumerate(row):
        x = 1.0 + ci * 2.8
        card_shape(s, x, y, 2.6, 0.75)
        is_loss = val.startswith("-")
        text_box(s, x + 0.1, y + 0.1, 2.4, 0.55, val, font_size=16,
                color=RGBColor(0xEF,0x44,0x44) if is_loss else (ACCENT_BLUE if ci > 0 else NAVY),
                bold=(ci==0), align=PP_ALIGN.CENTER)

# Milestones
text_box(s, 1.0, 6.0, 11, 0.5, "关键里程碑：第二年盈亏平衡 → 第三年首轮20万融资（出让10%股权）→ 第五年全感官助眠生态",
         font_size=14, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

add_page_number(s, 15)

# ========================================
# SLIDE 16: PRODUCT ROADMAP
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "timeline-bg.png"))
section_badge(s, "产品迭代")
sub_header(s, "全感官实体助眠生态——从压力到声光的多维干预")
# (guosai icon removed - using original design)

# Three pillars
pillars = [
    ("听觉", "白噪音枕头", "内置微型扬声器的智能枕头\n白噪音/自然音/脑波引导\n物理硬件，无需佩戴设备\n与加重被套形成配套生态"),
    ("视觉", "床上用品整体设计", "整体卧室视觉设计系统\n色彩心理学应用\n舒适视觉环境营造\n打造品牌专属美学"),
    ("触觉+嗅觉", "加重模块系统（现有）", "五区独立配重\n中草药缓释\n可调节精准深压\n双通道放松"),
]
for i, (sense, title, desc) in enumerate(pillars):
    x = 0.5 + i * 4.2
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3.8), Inches(5.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x40)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(2)
    # Sense badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(1.6), Inches(1.4), Inches(1.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sense
    run.font.size = Pt(18)
    run.font.color.rgb = NAVY
    run.font.bold = True
    run.font.name = "等线"
    # Title
    text_box(s, x + 0.3, 3.2, 3.2, 0.5, title, font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Description
    for j, line in enumerate(desc.split('\n')):
        text_box(s, x + 0.3, 3.9 + j * 0.5, 3.2, 0.4, line, font_size=13, color=LIGHT_GOLD, align=PP_ALIGN.CENTER)

# Bottom statement
text_box(s, 0.5, 6.6, 12.3, 0.5,
         "我们专注于有重量、可触达的实体硬件矩阵，构建属于品牌自身的硬件生态壁垒",
         font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_page_number(s, 16)

# ========================================
# SLIDE 17: TECH MOAT
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "技术护城河")
sub_header(s, "三位一体保护——支撑商业落地的技术壁垒")
# (guosai icon removed - using original design)

moat_items = [
    ("实用新型专利", "产品结构专利保护", "模块化加重被套的结构设计已获实用新型专利保护，防止直接抄袭。独特的分区设计与魔术贴固定系统构成核心技术壁垒。", "📜"),
    ("配方工艺专利", "中草药缓释技术", "薰衣草、柠檬草、柠檬香蜂草的配比与缓释工艺构成配方壁垒。体温激活、持续释放的技术方案难以逆向工程。", "⚗️"),
    ("商业数据壁垒", "用户数据与迭代", "积累用户睡眠偏好数据，持续优化模块配重方案。数据飞轮效应形成竞争壁垒，后来者难以复制。", "📊"),
]
for i, (title, subtitle, desc, icon) in enumerate(moat_items):
    x = 0.5 + i * 4.2
    card_shape(s, x, 1.3, 3.8, 5.5)
    # Icon
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(1.6), Inches(1.4), Inches(1.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
    circle.line.color.rgb = ACCENT_BLUE
    circle.line.width = Pt(2)
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.size = Pt(32)
    text_box(s, x + 0.3, 3.2, 3.2, 0.4, title, font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.3, 3.7, 3.2, 0.3, subtitle, font_size=13, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    gold_line(s, x + 0.8, 4.1, 2.2)
    text_box(s, x + 0.3, 4.3, 3.2, 2.2, desc, font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_page_number(s, 17)

# ========================================
# SLIDE 18: TEAM
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "团队介绍")
sub_header(s, "产学研深度融合——两位导师+八位核心成员")
# (guosai icon removed - using original design)

# Advisors
text_box(s, 0.5, 1.3, 12.3, 0.5, "指导教师", font_size=20, color=NAVY, bold=True)
gold_line(s, 0.5, 1.85, 2)

advisors = [("吴兆书", "指导教师"), ("葛峰", "指导教师")]
for i, (name, role) in enumerate(advisors):
    x = 0.5 + i * 3.0
    card_shape(s, x, 2.0, 2.7, 1.0)
    text_box(s, x + 0.2, 2.1, 2.3, 0.4, name, font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.2, 2.5, 2.3, 0.3, role, font_size=13, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# Team leader
text_box(s, 7.0, 1.3, 6.0, 0.5, "项目负责人", font_size=20, color=NAVY, bold=True)
gold_line(s, 7.0, 1.85, 2)
card_shape(s, 7.0, 2.0, 5.8, 1.0)
text_box(s, 7.3, 2.1, 5.2, 0.4, "高通  ·  项目负责人/CEO", font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 7.3, 2.5, 5.2, 0.3, "统筹项目全局，负责产品设计与商业落地", font_size=13, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# Team members
text_box(s, 0.5, 3.3, 12.3, 0.5, "核心成员", font_size=20, color=NAVY, bold=True)
gold_line(s, 0.5, 3.85, 2)

# Member roles
roles = [
    ("产品研发", "负责产品结构设计、材料选型、工艺优化"),
    ("文献研究", "负责医学文献检索、循证研究、学术支撑"),
    ("电商运营", "负责淘宝店铺运营、产品上架、客服管理"),
    ("直播运营", "负责抖音直播策划、短视频制作、流量获取"),
    ("市场营销", "负责品牌推广、用户调研、市场分析"),
    ("财务管理", "负责财务报表、成本核算、融资准备"),
    ("视觉设计", "负责产品包装、品牌VI、宣传物料"),
    ("技术支持", "负责小程序开发、数据分析、技术支持"),
]
for i, (role, desc) in enumerate(roles):
    row, col = divmod(i, 4)
    x = 0.5 + col * 3.15
    y = 4.1 + row * 1.5
    card_shape(s, x, y, 2.9, 1.3)
    text_box(s, x + 0.2, y + 0.1, 2.5, 0.4, role, font_size=16, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.2, y + 0.55, 2.5, 0.6, desc, font_size=11, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# Company registration note
text_box(s, 0.5, 6.8, 12.3, 0.4, "公司已注册完成，营业执照详见策划书附件",
         font_size=12, color=SUB_TEXT, align=PP_ALIGN.CENTER)

add_page_number(s, 18)

# ========================================
# SLIDE 19: SOCIAL IMPACT
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "content-bg.png"))
section_badge(s, "社会担当")
sub_header(s, "响应健康中国2030——关注青少年心理健康与睡眠保障")
# (guosai icon removed - using original design)

# Main message
text_box(s, 0.5, 1.5, 12.3, 0.8,
         "我们的项目不仅仅是一次创业尝试，更是一场极致的产学研大练兵",
         font_size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

gold_line(s, 4.0, 2.4, 5.3)

# Three achievements
achievements = [
    ("产学结合", "成功打通从顶层医学文献转化、产品结构开模验证，到全域电商落地的完整商业闭环"),
    ("健康中国", "响应《健康中国2030规划纲要》，关注青少年心理健康与睡眠保障，为青年人身心健康贡献力量"),
    ("教育成果", "团队在项目推进过程中，阅读文献、产学结合的能力大大提升，淘宝和抖音运营经历锻炼了实战能力"),
]
for i, (title, desc) in enumerate(achievements):
    x = 0.5 + i * 4.2
    card_shape(s, x, 2.8, 3.8, 3.5)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.4), Inches(3.0), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(i+1)
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.bold = True
    text_box(s, x + 0.3, 4.2, 3.2, 0.4, title, font_size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.3, 4.7, 3.2, 1.4, desc, font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_page_number(s, 19)

# ========================================
# SLIDE 20: THANK YOU
# ========================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, os.path.join(BG, "thankyou-bg.png"))

# Gold border decoration
add_image_safe(s, os.path.join(BG, "gold-line.png"), 4.0, 1.5, width=5.3)

# Thank you text
text_box(s, 0, 2.0, 13.33, 1.0, "感谢聆听", font_size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 0, 3.2, 13.33, 0.6, "助你好梦——重新定义现代压力下的深度睡眠方案", font_size=20, color=LIGHT_GOLD, align=PP_ALIGN.CENTER)

gold_line(s, 4.5, 4.0, 4.3)

# Contact info
text_box(s, 0, 4.5, 13.33, 0.5, "助你好梦创业团队", font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
text_box(s, 0, 5.2, 13.33, 0.5, "中国国际大学生创新大赛（2025）", font_size=16, color=LIGHT_GOLD, align=PP_ALIGN.CENTER)